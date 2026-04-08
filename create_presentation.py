#!/usr/bin/env python3
"""
Gera a apresentação PowerPoint do projeto EDA — Olist.
Uso: python3 create_presentation.py
"""

import os
import pandas as pd
import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import seaborn as sns
from sklearn.model_selection import train_test_split
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

# ================================================================
# CONFIG
# ================================================================
DATA_PATH  = "data/"
ASSETS_DIR = "slides_assets"
OUTPUT     = "apresentacao_olist.pptx"
os.makedirs(ASSETS_DIR, exist_ok=True)

C_NAVY  = RGBColor(0x1E, 0x29, 0x3B)
C_BLUE  = RGBColor(0x3B, 0x82, 0xF6)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_LGRAY = RGBColor(0xF1, 0xF5, 0xF9)
C_DGRAY = RGBColor(0x64, 0x74, 0x8B)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

SCORE_COLORS = ["#EF4444", "#F97316", "#EAB308", "#84CC16", "#22C55E"]

# ================================================================
# 1. CARREGAR E PREPARAR DADOS
# ================================================================
print("Carregando dados...")

orders     = pd.read_csv(DATA_PATH + "olist_orders_dataset.csv")
customers  = pd.read_csv(DATA_PATH + "olist_customers_dataset.csv")
items      = pd.read_csv(DATA_PATH + "olist_order_items_dataset.csv")
payments   = pd.read_csv(DATA_PATH + "olist_order_payments_dataset.csv")
reviews    = pd.read_csv(DATA_PATH + "olist_order_reviews_dataset.csv")
products   = pd.read_csv(DATA_PATH + "olist_products_dataset.csv")
categories = pd.read_csv(DATA_PATH + "product_category_name_translation.csv")

items_agg = items.groupby("order_id").agg(
    n_items=("order_item_id", "count"),
    total_price=("price", "sum"),
    total_freight=("freight_value", "sum")
).reset_index()

payments_agg = payments.groupby("order_id").agg(
    payment_value=("payment_value", "sum"),
    payment_type=("payment_type", lambda x: x.mode()[0]),
    payment_installments=("payment_installments", "max")
).reset_index()

reviews_target = reviews.groupby("order_id")["review_score"].first().reset_index()

df = (
    orders
    .merge(customers[["customer_id", "customer_state"]], on="customer_id", how="left")
    .merge(items_agg, on="order_id", how="left")
    .merge(payments_agg, on="order_id", how="left")
    .merge(reviews_target, on="order_id", how="left")
)

for col in ["order_purchase_timestamp", "order_delivered_customer_date", "order_estimated_delivery_date"]:
    df[col] = pd.to_datetime(df[col])

df["delivery_time_days"] = (
    df["order_delivered_customer_date"] - df["order_purchase_timestamp"]
).dt.days

df["delivery_late"] = (
    (df["order_delivered_customer_date"] > df["order_estimated_delivery_date"]) &
    df["order_delivered_customer_date"].notna()
).astype(int)

df = df.dropna(subset=["review_score"])
df["review_score"] = df["review_score"].astype(int)

drop_cols = [
    "order_id", "customer_id", "order_purchase_timestamp", "order_approved_at",
    "order_delivered_carrier_date", "order_delivered_customer_date",
    "order_estimated_delivery_date", "review_score"
]
X = df.drop(columns=drop_cols)
y = df["review_score"]

X_train, X_test, y_train, y_test = train_test_split(
    X, y, test_size=0.2, random_state=42, stratify=y
)
train_df = X_train.copy()
train_df["review_score"] = y_train.values
print(f"Dados prontos — X_train: {X_train.shape}")

# ================================================================
# 2. GERAR GRÁFICOS
# ================================================================
print("Gerando graficos...")
sns.set_theme(style="whitegrid", font_scale=1.1)

def save_fig(name):
    path = os.path.join(ASSETS_DIR, f"{name}.png")
    plt.savefig(path, dpi=150, bbox_inches="tight", facecolor="white")
    plt.close("all")
    return path

# --- Plot 1: Distribuição do target ---
score_counts = y_train.value_counts().sort_index()
fig, ax = plt.subplots(figsize=(7, 4))
bars = ax.bar(score_counts.index, score_counts.values,
              color=SCORE_COLORS, edgecolor="white", linewidth=0.5, zorder=3)
for bar, val in zip(bars, score_counts.values):
    pct = val / len(y_train) * 100
    ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + 150,
            f"{val:,}\n({pct:.1f}%)", ha="center", va="bottom", fontsize=9, color="#374151")
ax.set_xticks([1, 2, 3, 4, 5])
ax.set_xlabel("Nota de Avaliação (estrelas)", labelpad=8)
ax.set_ylabel("Nº de Pedidos")
ax.set_title("Distribuição do Target — review_score", fontsize=13, pad=10)
ax.grid(axis="y", alpha=0.4, zorder=0)
ax.set_axisbelow(True)
plt.tight_layout()
plot_target = save_fig("01_target")

# --- Plot 2: delivery_late vs score ---
fig, ax = plt.subplots(figsize=(5.5, 4))
sns.boxplot(data=train_df, x="delivery_late", y="review_score",
            palette=["#22C55E", "#EF4444"], ax=ax, width=0.45, linewidth=1.2)
ax.set_xticklabels(["No prazo", "Atrasado"], fontsize=12)
ax.set_xlabel("")
ax.set_ylabel("Nota de Avaliação")
ax.set_yticks([1, 2, 3, 4, 5])
ax.set_title("Atraso na Entrega vs Nota", fontsize=13, pad=10)
plt.tight_layout()
plot_late = save_fig("02_delivery_late")

# --- Plot 3: delivery_time_days vs score ---
cap = train_df["delivery_time_days"].quantile(0.97)
plot_data3 = train_df[train_df["delivery_time_days"] <= cap].dropna(subset=["delivery_time_days"])
fig, ax = plt.subplots(figsize=(8, 4))
sns.boxplot(data=plot_data3, x="review_score", y="delivery_time_days",
            palette=SCORE_COLORS, ax=ax, linewidth=1.2)
ax.set_xlabel("Nota de Avaliação")
ax.set_ylabel("Dias até a Entrega")
ax.set_title("Tempo de Entrega (dias) por Nota", fontsize=13, pad=10)
plt.tight_layout()
plot_time = save_fig("03_delivery_time")

# --- Plot 4: pedidos por mês ---
train_orders = df.loc[X_train.index, ["order_purchase_timestamp"]].copy()
train_orders["month"] = train_orders["order_purchase_timestamp"].dt.to_period("M")
monthly = train_orders["month"].value_counts().sort_index()
fig, ax = plt.subplots(figsize=(11, 3.5))
monthly.plot(kind="bar", color="#3B82F6", edgecolor="white", ax=ax, zorder=3)
ax.set_ylabel("Nº de Pedidos")
ax.set_title("Volume de Pedidos por Mês", fontsize=13, pad=10)
ax.grid(axis="y", alpha=0.4, zorder=0)
ax.set_axisbelow(True)
plt.xticks(rotation=45, ha="right")
plt.tight_layout()
plot_monthly = save_fig("04_monthly")

# --- Plot 5: top 10 categorias ---
train_order_ids = df.loc[X_train.index, "order_id"]
tic = (
    items
    .merge(products[["product_id", "product_category_name"]], on="product_id", how="left")
    .merge(categories, on="product_category_name", how="left")
)
tic = tic[tic["order_id"].isin(train_order_ids)]
top_cats = tic["product_category_name_english"].dropna().value_counts().head(10)
fig, ax = plt.subplots(figsize=(8, 4.5))
colors_bars = ["#3B82F6" if i == 0 else "#93C5FD" for i in range(10)]
ax.barh(top_cats.index[::-1], top_cats.values[::-1],
        color=colors_bars[::-1], edgecolor="white")
ax.set_xlabel("Nº de Itens Vendidos")
ax.set_title("Top 10 Categorias de Produto", fontsize=13, pad=10)
ax.grid(axis="x", alpha=0.4)
ax.set_axisbelow(True)
plt.tight_layout()
plot_cats = save_fig("05_categories")

# --- Plot 6: heatmap de correlação ---
corr_features = ["total_price", "total_freight", "payment_value",
                 "payment_installments", "n_items", "delivery_time_days", "delivery_late"]
corr_df2 = X_train[corr_features].copy()
corr_df2["review_score"] = y_train.values
fig, ax = plt.subplots(figsize=(8, 6))
sns.heatmap(corr_df2.corr(), annot=True, fmt=".2f", cmap="RdBu_r",
            center=0, square=True, ax=ax, linewidths=0.5, annot_kws={"size": 9})
ax.set_title("Matriz de Correlação", fontsize=13, pad=10)
plt.tight_layout()
plot_corr = save_fig("06_correlation")

print("Graficos salvos.")

# ================================================================
# 3. HELPERS PARA CONSTRUIR O PPTX
# ================================================================

def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, left, top, w, h, color):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, w, h
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def textbox(slide, text, left, top, w, h,
            size=16, bold=False, italic=False,
            color=None, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(left, top, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color or C_NAVY
    return txb

def header(slide, title):
    rect(slide, 0, 0, SLIDE_W, Inches(1.1), C_NAVY)
    textbox(slide, title,
            Inches(0.55), Inches(0.15),
            Inches(12.2), Inches(0.85),
            size=24, bold=True, color=C_WHITE)

def bullets(slide, items_list, left, top, w, h, size=16, color=None):
    txb = slide.shapes.add_textbox(left, top, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items_list):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(5)
        run = p.add_run()
        run.text = f"▸  {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = color or C_NAVY

def pic(slide, path, left, top, w, h):
    slide.shapes.add_picture(path, left, top, w, h)

# ================================================================
# 4. MONTAR OS SLIDES
# ================================================================
print("Montando slides...")
prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

# ── Slide 1: Título ─────────────────────────────────────────────
s = new_slide(prs)
set_bg(s, C_NAVY)
rect(s, Inches(0.55), Inches(3.0), Inches(1.4), Inches(0.07), C_BLUE)
textbox(s, "Previsão de Satisfação do Cliente",
        Inches(0.55), Inches(1.4), Inches(12.2), Inches(1.4),
        size=40, bold=True, color=C_WHITE)
textbox(s, "Análise Exploratória de Dados — Brazilian E-Commerce (Olist)",
        Inches(0.55), Inches(3.2), Inches(12.2), Inches(0.8),
        size=20, color=RGBColor(0xAB, 0xB8, 0xC3))
textbox(s, "Mineração de Dados Aplicada  ·  CIn / UFPE  ·  2025",
        Inches(0.55), Inches(6.7), Inches(12.2), Inches(0.5),
        size=12, color=RGBColor(0x64, 0x74, 0x8B))

# ── Slide 2: Agenda ─────────────────────────────────────────────
s = new_slide(prs)
set_bg(s, C_WHITE)
header(s, "Agenda")

agenda = [
    ("1", "Contexto e domínio"),
    ("2", "Problema de negócio & variável-alvo"),
    ("3", "Estrutura do dataset"),
    ("4", "Metodologia (train/test split)"),
    ("5", "Análise exploratória — principais achados"),
    ("6", "Conclusões e próximos passos"),
]
cols = [Inches(0.6), Inches(7.0)]
for i, (num, label) in enumerate(agenda):
    col = cols[0] if i < 3 else cols[1]
    row = i if i < 3 else i - 3
    y   = Inches(1.4) + row * Inches(1.8)
    rect(s, col, y, Inches(0.65), Inches(0.65), C_BLUE)
    textbox(s, num, col, y, Inches(0.65), Inches(0.65),
            size=18, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    textbox(s, label, col + Inches(0.75), y + Inches(0.1),
            Inches(5.5), Inches(0.55), size=17, color=C_NAVY)

# ── Slide 3: Contexto ───────────────────────────────────────────
s = new_slide(prs)
set_bg(s, C_WHITE)
header(s, "Contexto e Domínio")

bullets(s, [
    "Dataset público da Olist disponibilizado no Kaggle",
    "Dados reais e anonimizados de marketplaces brasileiros",
    "Período coberto: 2016 a 2018",
    "~100 mil pedidos distribuídos em 8 arquivos CSV",
    "Perspectivas: produtos, pagamentos, logística, clientes e avaliações",
], Inches(0.6), Inches(1.3), Inches(7.8), Inches(5.5), size=18)

stats = [("99.4 mil", "Pedidos"), ("96.1 mil", "Clientes únicos"),
         ("32.9 mil", "Produtos"), ("3.1 mil", "Vendedores")]
for i, (val, label) in enumerate(stats):
    bx, by = Inches(9.5), Inches(1.3) + i * Inches(1.5)
    rect(s, bx, by, Inches(3.3), Inches(1.2), C_BLUE)
    textbox(s, val, bx, by + Inches(0.05), Inches(3.3), Inches(0.65),
            size=24, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    textbox(s, label, bx, by + Inches(0.65), Inches(3.3), Inches(0.45),
            size=13, color=C_WHITE, align=PP_ALIGN.CENTER)

# ── Slide 4: Problema de Negócio ────────────────────────────────
s = new_slide(prs)
set_bg(s, C_WHITE)
header(s, "Problema de Negócio")

rect(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(1.5), C_LGRAY)
textbox(s,
        '"Prever a nota de satisfação do cliente após uma compra,\n'
        'identificando pedidos com maior risco de gerar avaliações ruins."',
        Inches(0.7), Inches(1.4), Inches(12.0), Inches(1.3),
        size=18, italic=True, color=C_NAVY)

textbox(s, "Por que isso importa?",
        Inches(0.5), Inches(3.05), Inches(12.3), Inches(0.5),
        size=17, bold=True, color=C_BLUE)
bullets(s, [
    "Avaliações negativas afetam confiança, recompra e reputação do marketplace",
    "Permite ação preventiva: contato proativo, priorização de suporte, revisão logística",
    "Detectar pedidos com risco de 1-2 estrelas antes da entrega gera valor real",
], Inches(0.6), Inches(3.55), Inches(12.0), Inches(2.0), size=16)

rect(s, Inches(0.5), Inches(5.9), Inches(12.3), Inches(1.2), C_NAVY)
textbox(s, "Variável-alvo:  review_score  (1 a 5 estrelas)  —  Classificação Multiclasse",
        Inches(0.7), Inches(6.05), Inches(7.5), Inches(0.6),
        size=16, bold=True, color=C_WHITE)
textbox(s, "⚠  Atenção ética: não responsabilizar\nexclusivamente vendedores por falhas logísticas.",
        Inches(8.5), Inches(5.95), Inches(4.2), Inches(0.9),
        size=12, color=RGBColor(0xAB, 0xB8, 0xC3))

# ── Slide 5: Dataset ────────────────────────────────────────────
s = new_slide(prs)
set_bg(s, C_WHITE)
header(s, "Estrutura do Dataset")

rows = [
    ("orders",          "99.441 × 8",   "Pedidos: status, datas de compra, aprovação e entrega"),
    ("customers",       "99.441 × 5",   "Clientes: ID único, cidade, estado"),
    ("order_items",     "112.650 × 7",  "Itens: produto, vendedor, preço, frete"),
    ("order_payments",  "103.886 × 5",  "Pagamentos: tipo, parcelas, valor"),
    ("order_reviews",   "99.224 × 7",   "Avaliações: nota, comentário, data  ← TARGET"),
    ("products",        "32.951 × 9",   "Produtos: categoria, dimensões, peso"),
    ("sellers",         "3.095 × 4",    "Vendedores: localização"),
    ("categories",      "71 × 2",       "Tradução de categorias PT → EN"),
]
col_xs = [Inches(0.4), Inches(3.3), Inches(5.3)]
col_ws = [Inches(2.8), Inches(1.9), Inches(7.6)]
hdrs   = ["Arquivo", "Dimensão", "Conteúdo"]
row_h  = Inches(0.57)

for ci, (hdr_txt, cx, cw) in enumerate(zip(hdrs, col_xs, col_ws)):
    rect(s, cx, Inches(1.2), cw, row_h, C_NAVY)
    textbox(s, hdr_txt, cx + Inches(0.07), Inches(1.2),
            cw, row_h, size=13, bold=True, color=C_WHITE)

for ri, (fname, dim, desc) in enumerate(rows):
    y   = Inches(1.2) + row_h + ri * row_h
    bg  = C_LGRAY if ri % 2 == 0 else C_WHITE
    hl  = "reviews" in fname
    fc  = C_BLUE if hl else C_NAVY
    for val, cx, cw in zip([fname, dim, desc], col_xs, col_ws):
        rect(s, cx, y, cw, row_h, bg)
        textbox(s, val, cx + Inches(0.07), y + Inches(0.07),
                cw - Inches(0.1), row_h, size=12, bold=hl, color=fc)

# ── Slide 6: Train/Test Split ───────────────────────────────────
s = new_slide(prs)
set_bg(s, C_WHITE)
header(s, "Metodologia — Train/Test Split")

textbox(s, "A Regra de Ouro: separar ANTES de qualquer análise",
        Inches(0.5), Inches(1.25), Inches(12.3), Inches(0.5),
        size=18, bold=True, color=C_NAVY)
textbox(s,
        "Toda a EDA foi feita exclusivamente com os dados de treino (X_train, y_train).\n"
        "Isso evita o Data Leakage — vazamento de informação do teste para o treino.",
        Inches(0.5), Inches(1.85), Inches(12.3), Inches(0.9),
        size=15, color=C_DGRAY)

rect(s, Inches(1.5), Inches(3.05), Inches(10.3), Inches(0.75), C_NAVY)
textbox(s, "Dataset completo  (99.224 pedidos com avaliação)",
        Inches(1.5), Inches(3.05), Inches(10.3), Inches(0.75),
        size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
textbox(s, "▼  train_test_split(test_size=0.2, stratify=y, random_state=42)",
        Inches(3.5), Inches(3.9), Inches(6.3), Inches(0.4),
        size=13, color=C_DGRAY, align=PP_ALIGN.CENTER)

rect(s, Inches(0.8), Inches(4.4), Inches(6.1), Inches(1.2), C_BLUE)
textbox(s, "X_train  +  y_train\n80%  —  ~79.4 mil pedidos\nUsado na EDA e no treino do modelo",
        Inches(0.8), Inches(4.4), Inches(6.1), Inches(1.2),
        size=14, color=C_WHITE, align=PP_ALIGN.CENTER)

rect(s, Inches(7.4), Inches(4.4), Inches(5.1), Inches(1.2), C_DGRAY)
textbox(s, "X_test  +  y_test\n20%  —  ~19.8 mil pedidos\nSó usado na avaliação final",
        Inches(7.4), Inches(4.4), Inches(5.1), Inches(1.2),
        size=14, color=C_WHITE, align=PP_ALIGN.CENTER)

textbox(s,
        "stratify=y  garante que a proporção de cada nota (1–5) seja preservada nos dois conjuntos.",
        Inches(0.5), Inches(5.8), Inches(12.3), Inches(0.5),
        size=13, italic=True, color=C_DGRAY, align=PP_ALIGN.CENTER)

# ── Slide 7: Distribuição do Target ─────────────────────────────
s = new_slide(prs)
set_bg(s, C_WHITE)
header(s, "Distribuição da Variável-Alvo")

pic(s, plot_target, Inches(4.8), Inches(1.2), Inches(8.0), Inches(5.2))
bullets(s, [
    "Nota 5 representa ~57% dos pedidos",
    "Classes 2 e 3 são as mais raras",
    "Dataset fortemente desbalanceado",
    "Acurácia simples será enganosa",
    "Métrica correta:\nweighted F1  ou  macro F1",
], Inches(0.4), Inches(1.5), Inches(4.1), Inches(5.0), size=16)

# ── Slide 8: Logística vs Satisfação ────────────────────────────
s = new_slide(prs)
set_bg(s, C_WHITE)
header(s, "Principal Achado — Impacto da Logística na Satisfação")

pic(s, plot_late, Inches(0.3), Inches(1.2), Inches(5.4), Inches(3.9))
pic(s, plot_time, Inches(5.9), Inches(1.2), Inches(7.2), Inches(3.9))

rect(s, Inches(0.3), Inches(5.25), Inches(12.8), Inches(1.7), C_LGRAY)
textbox(s,
        "Conclusao: pedidos atrasados recebem notas significativamente menores.\n"
        "Quanto maior o tempo de entrega, menor a nota.\n"
        "A logística é o principal driver de insatisfação neste dataset.",
        Inches(0.5), Inches(5.35), Inches(12.5), Inches(1.5),
        size=15, bold=False, color=C_NAVY)

# ── Slide 9: Categorias e Temporal ──────────────────────────────
s = new_slide(prs)
set_bg(s, C_WHITE)
header(s, "Categorias de Produto e Comportamento Temporal")

pic(s, plot_cats,    Inches(0.3), Inches(1.2), Inches(6.3), Inches(5.2))
pic(s, plot_monthly, Inches(6.8), Inches(1.2), Inches(6.2), Inches(5.2))

textbox(s, "Top 10 categorias por volume de itens",
        Inches(0.3), Inches(6.5), Inches(6.3), Inches(0.4),
        size=11, italic=True, color=C_DGRAY, align=PP_ALIGN.CENTER)
textbox(s, "Crescimento acelerado em 2017-2018 — pico em nov. (Black Friday)",
        Inches(6.8), Inches(6.5), Inches(6.2), Inches(0.4),
        size=11, italic=True, color=C_DGRAY, align=PP_ALIGN.CENTER)

# ── Slide 10: Correlações ────────────────────────────────────────
s = new_slide(prs)
set_bg(s, C_WHITE)
header(s, "Correlações entre Variáveis Numéricas")

pic(s, plot_corr, Inches(0.3), Inches(1.2), Inches(7.2), Inches(5.5))
bullets(s, [
    "delivery_late e review_score: correlação negativa — confirma impacto do atraso",
    "delivery_time_days e review_score: correlação negativa — quanto mais demora, pior a nota",
    "total_price e payment_value: altamente correlacionados — candidato a remoção (redundância)",
    "n_items tem baixa correlação com o target",
    "payment_installments tem leve correlação positiva com o valor pago",
], Inches(7.8), Inches(1.5), Inches(5.3), Inches(5.2), size=14)

# ── Slide 11: Conclusões ─────────────────────────────────────────
s = new_slide(prs)
set_bg(s, C_WHITE)
header(s, "Conclusões da EDA")

conclusions = [
    ("Desbalanceamento", "Nota 5 domina (~57%). Usar weighted / macro F1 como métrica de avaliação."),
    ("Driver principal", "Logística (atraso + tempo de entrega) é o fator mais correlacionado com insatisfação."),
    ("Valores ausentes", "delivery_time_days tem NaNs em pedidos não entregues — imputar com mediana ou criar flag binário."),
    ("Variação regional", "Nota média varia por estado — indicativo de qualidade logística regional desigual."),
    ("Sazonalidade", "Crescimento acelerado em 2017-2018; pico em novembro (Black Friday). Sazonalidade relevante."),
]
for i, (title, desc) in enumerate(conclusions):
    y_pos = Inches(1.3) + i * Inches(1.1)
    rect(s, Inches(0.4), y_pos + Inches(0.08), Inches(0.07), Inches(0.75), C_BLUE)
    textbox(s, title, Inches(0.62), y_pos, Inches(3.0), Inches(0.45),
            size=15, bold=True, color=C_NAVY)
    textbox(s, desc, Inches(0.62), y_pos + Inches(0.44), Inches(12.3), Inches(0.55),
            size=13, color=C_DGRAY)

# ── Slide 12: Próximos Passos ────────────────────────────────────
s = new_slide(prs)
set_bg(s, C_NAVY)

textbox(s, "Próximos Passos",
        Inches(0.55), Inches(0.5), Inches(12.3), Inches(0.8),
        size=32, bold=True, color=C_WHITE)
rect(s, Inches(0.55), Inches(1.35), Inches(1.2), Inches(0.07), C_BLUE)

next_steps = [
    ("1", "Pré-processamento",    "Tratar NaNs, encoding de categóricas, normalização"),
    ("2", "Feature Engineering",  "Novas features derivadas de logística, temporais e de produto"),
    ("3", "Modelagem",            "Testar Random Forest, XGBoost e outros classificadores multiclasse"),
    ("4", "Avaliação",            "Macro F1 e weighted F1 como métricas principais; análise por classe"),
    ("5", "Interpretabilidade",   "SHAP values para explicar predições e monitorar viés por região/categoria"),
]
cols_ns = [Inches(0.55), Inches(7.1)]
for i, (num, title, desc) in enumerate(next_steps):
    cx = cols_ns[0] if i < 3 else cols_ns[1]
    row = i if i < 3 else i - 3
    y  = Inches(1.6) + row * Inches(1.8)
    rect(s, cx, y, Inches(0.7), Inches(0.7), C_BLUE)
    textbox(s, num, cx, y, Inches(0.7), Inches(0.7),
            size=22, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    textbox(s, title, cx + Inches(0.8), y, Inches(5.7), Inches(0.45),
            size=17, bold=True, color=C_WHITE)
    textbox(s, desc, cx + Inches(0.8), y + Inches(0.45), Inches(5.7), Inches(0.6),
            size=13, color=RGBColor(0xAB, 0xB8, 0xC3))

# ================================================================
# 5. SALVAR
# ================================================================
prs.save(OUTPUT)
print(f"Apresentacao salva em: {OUTPUT}")
